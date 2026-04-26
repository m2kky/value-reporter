"""Premium PPTX generator with ValueWats brand identity.

Generates agency-quality presentations with:
- Custom dark theme with lime accent (#e2f300)
- Logo on cover and closing slides
- Gradient section dividers
- Formatted KPI cards and data tables
"""
from __future__ import annotations

import json
import re
from pathlib import Path
import logging
from typing import Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)

# ── Brand Palette ──────────────────────────────────────────
LIME       = RGBColor(0xE2, 0xF3, 0x00) if PPTX_AVAILABLE else None
BG_DARK    = RGBColor(0x23, 0x23, 0x18) if PPTX_AVAILABLE else None
TEXT_CREAM = RGBColor(0xFF, 0xFE, 0xD9) if PPTX_AVAILABLE else None
OLIVE      = RGBColor(0x7A, 0x78, 0x39) if PPTX_AVAILABLE else None
MUTED      = RGBColor(0xB7, 0xB5, 0x7E) if PPTX_AVAILABLE else None
CARD_BG    = RGBColor(0x2A, 0x2A, 0x1F) if PPTX_AVAILABLE else None
LOGO_BLUE  = RGBColor(0x28, 0x59, 0xAF) if PPTX_AVAILABLE else None
WHITE      = RGBColor(0xFF, 0xFF, 0xFF) if PPTX_AVAILABLE else None
RED_ACCENT = RGBColor(0xEF, 0x44, 0x44) if PPTX_AVAILABLE else None
GREEN_ACC  = RGBColor(0x10, 0xB9, 0x81) if PPTX_AVAILABLE else None

LOGO_PATH = Path(__file__).resolve().parent.parent / "value logo.png"

SLIDE_WIDTH  = Inches(13.333) if PPTX_AVAILABLE else 0
SLIDE_HEIGHT = Inches(7.5)    if PPTX_AVAILABLE else 0


def _set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color=None):
    """Add a rectangle shape to a slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()  # No border
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  color=None, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or TEXT_CREAM
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size=14,
                     color=None, bullet_color=None):
    """Add a bulleted list with styled bullets."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or TEXT_CREAM
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox


def _cover_slide(prs, client_name: str, period: str):
    """Create a branded cover slide with logo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide, BG_DARK)

    # Lime accent bar at top
    _add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), LIME)

    # Logo (centered, upper area)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH), Inches(4.5), Inches(1.2), width=Inches(4.3)
        )

    # Client name
    _add_text_box(
        slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.9),
        client_name, font_size=36, color=TEXT_CREAM, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Report title
    _add_text_box(
        slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.6),
        "Monthly Performance Report", font_size=22, color=MUTED,
        alignment=PP_ALIGN.CENTER,
    )

    # Period
    _add_text_box(
        slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5),
        period, font_size=18, color=LIME, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Bottom accent bar
    _add_shape(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), LIME)


def _section_divider(prs, title: str, subtitle: str = ""):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)

    # Left accent bar
    _add_shape(slide, Inches(0.8), Inches(2.5), Inches(0.08), Inches(2.5), LIME)

    _add_text_box(
        slide, Inches(1.2), Inches(2.8), Inches(10), Inches(1),
        title, font_size=40, color=TEXT_CREAM, bold=True,
    )
    if subtitle:
        _add_text_box(
            slide, Inches(1.2), Inches(3.8), Inches(10), Inches(0.6),
            subtitle, font_size=18, color=MUTED,
        )


def _kpi_card(slide, left, top, label: str, value: str, change: str = ""):
    """Add a single KPI card."""
    card_w, card_h = Inches(2.8), Inches(1.6)
    _add_shape(slide, left, top, card_w, card_h, CARD_BG)

    # Lime top edge
    _add_shape(slide, left, top, card_w, Inches(0.05), LIME)

    _add_text_box(
        slide, left + Inches(0.2), top + Inches(0.15), card_w - Inches(0.4), Inches(0.4),
        label, font_size=11, color=MUTED,
    )
    _add_text_box(
        slide, left + Inches(0.2), top + Inches(0.55), card_w - Inches(0.4), Inches(0.6),
        value, font_size=28, color=TEXT_CREAM, bold=True,
    )
    if change:
        is_positive = change.startswith("+")
        is_negative = change.startswith("-")
        change_color = GREEN_ACC if is_positive else (RED_ACCENT if is_negative else MUTED)
        _add_text_box(
            slide, left + Inches(0.2), top + Inches(1.15), card_w - Inches(0.4), Inches(0.3),
            change, font_size=12, color=change_color, bold=True,
        )


def _kpi_slide(prs, title: str, kpis: list[tuple[str, str, str]]):
    """Create a KPI dashboard slide with up to 8 cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)

    _add_text_box(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
        title, font_size=28, color=TEXT_CREAM, bold=True,
    )
    _add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), LIME)

    cols = 4
    for i, (label, value, change) in enumerate(kpis[:8]):
        row = i // cols
        col = i % cols
        left = Inches(0.5) + col * Inches(3.1)
        top = Inches(1.4) + row * Inches(2.0)
        _kpi_card(slide, left, top, label, value, change)


def _data_table_slide(prs, title: str, headers: list[str],
                      rows: list[list[str]], max_rows: int = 8):
    """Create a slide with a formatted data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)

    _add_text_box(
        slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
        title, font_size=24, color=TEXT_CREAM, bold=True,
    )
    _add_shape(slide, Inches(0.6), Inches(0.85), Inches(1.5), Inches(0.04), LIME)

    display_rows = rows[:max_rows]
    n_rows = len(display_rows) + 1  # +1 for header
    n_cols = len(headers)

    table_left = Inches(0.4)
    table_top = Inches(1.1)
    table_width = Inches(12.5)
    table_height = Inches(0.4) * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
    table = table_shape.table

    # Style header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = BG_DARK
            paragraph.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIME

    # Style data rows
    for i, row_data in enumerate(display_rows):
        bg = CARD_BG if i % 2 == 0 else BG_DARK
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = TEXT_CREAM
                paragraph.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg


def _bullets_slide(prs, title: str, items: list[str], subtitle: str = ""):
    """Create a slide with numbered/bulleted items."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)

    _add_text_box(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
        title, font_size=28, color=TEXT_CREAM, bold=True,
    )
    _add_shape(slide, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04), LIME)

    if subtitle:
        _add_text_box(
            slide, Inches(0.6), Inches(1.2), Inches(11), Inches(0.4),
            subtitle, font_size=14, color=MUTED,
        )

    numbered = []
    for i, item in enumerate(items[:10], 1):
        clean = re.sub(r'^\d+\.\s*', '', item).strip()
        clean = re.sub(r'^[-•]\s*', '', clean).strip()
        numbered.append(f"{i}.  {clean}")

    _add_bullet_list(
        slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5),
        numbered, font_size=15, color=TEXT_CREAM,
    )


def _closing_slide(prs):
    """Create a branded closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)

    _add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), LIME)

    _add_text_box(
        slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1),
        "Thank You", font_size=48, color=TEXT_CREAM, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
        "Powered by Value Marketing Solutions ©",
        font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER,
    )

    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH), Inches(5.2), Inches(4.5), width=Inches(3)
        )

    _add_shape(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), LIME)


# ── Markdown Parsing Helpers ────────────────────────────────

def _extract_section(md: str, heading: str) -> str:
    """Extract content under a specific heading from markdown."""
    pattern = rf"^##\s+{re.escape(heading)}\s*$(.*?)(?=^##\s|\Z)"
    match = re.search(pattern, md, re.MULTILINE | re.DOTALL)
    return match.group(1).strip() if match else ""


def _extract_bullets(text: str) -> list[str]:
    """Extract bullet points from markdown text."""
    return [line.strip().lstrip("-•").strip() for line in text.splitlines()
            if line.strip().startswith(("-", "•", "1", "2", "3", "4", "5"))]


def _extract_table(text: str) -> tuple[list[str], list[list[str]]]:
    """Extract a markdown table into headers and rows."""
    lines = [l.strip() for l in text.splitlines() if l.strip().startswith("|")]
    if len(lines) < 2:
        return [], []
    headers = [c.strip() for c in lines[0].split("|") if c.strip() and c.strip() != "---"]
    rows = []
    for line in lines[2:]:  # Skip header + separator
        cells = [c.strip() for c in line.split("|") if c.strip()]
        if cells:
            rows.append(cells)
    return headers, rows


def _pct_str(val: float | None) -> str:
    if val is None:
        return "n/a"
    sign = "+" if val > 0 else ""
    return f"{sign}{val:.1f}%"


def _money(val: float, cur: str) -> str:
    return f"{val:,.2f} {cur}"


def _num(val: float) -> str:
    return f"{val:,.0f}"


# ── Main Generator ──────────────────────────────────────────

def generate_pptx_from_data(
    output_path: Path,
    *,
    client_name: str,
    period: str,
    currency: str,
    current_kpis: dict[str, float],
    previous_kpis: dict[str, float],
    changes: dict[str, float | None],
    campaigns: list[dict[str, Any]],
    organic_summary: dict[str, Any] | None = None,
    conversations: dict[str, Any] | None = None,
) -> bool:
    """Generate a premium branded PPTX from structured data."""
    if not PPTX_AVAILABLE:
        logger.warning("python-pptx is not installed. Skipping PPTX generation.")
        return False

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ── 1. Cover ──
    _cover_slide(prs, client_name, period)

    # ── 2. Executive KPIs ──
    kpis = [
        ("Total Spend", _money(current_kpis.get("spend", 0), currency), _pct_str(changes.get("spend"))),
        ("Impressions", _num(current_kpis.get("impressions", 0)), _pct_str(changes.get("impressions"))),
        ("Reach", _num(current_kpis.get("reach", 0)), _pct_str(changes.get("reach"))),
        ("Clicks", _num(current_kpis.get("clicks", 0)), _pct_str(changes.get("clicks"))),
        ("CTR", f"{current_kpis.get('ctr', 0):.2f}%", _pct_str(changes.get("ctr"))),
        ("CPC", _money(current_kpis.get("cpc", 0), currency), _pct_str(changes.get("cpc"))),
        ("CPM", _money(current_kpis.get("cpm", 0), currency), _pct_str(changes.get("cpm"))),
        ("ROAS", f"{current_kpis.get('roas_meta', 0):.2f}", _pct_str(changes.get("roas_meta"))),
    ]
    _kpi_slide(prs, "Paid Performance Overview", kpis)

    # ── 3. Paid Section ──
    _section_divider(prs, "Paid Advertising", f"{len(campaigns)} campaigns analyzed")

    # Campaign table
    if campaigns:
        headers = ["Campaign", "Spend", "Clicks", "CTR", "Purchases", "ROAS"]
        rows = []
        for c in sorted(campaigns, key=lambda x: x.get("spend", 0), reverse=True)[:8]:
            name = str(c.get("campaign_name", "Unknown"))[:40]
            rows.append([
                name,
                _money(c.get("spend", 0), currency),
                _num(c.get("clicks", 0)),
                f"{c.get('ctr', 0):.2f}%",
                _num(c.get("purchases", 0)),
                f"{c.get('roas_meta', 0):.2f}",
            ])
        _data_table_slide(prs, "Top Campaigns by Spend", headers, rows)

    # ── 4. Organic Section ──
    if organic_summary and organic_summary.get("totals", {}).get("posts", 0):
        totals = organic_summary.get("totals", {})
        _section_divider(prs, "Organic Content", f"{int(totals.get('posts', 0))} pieces analyzed")

        org_kpis = [
            ("Content Published", _num(totals.get("posts", 0)), ""),
            ("Views", _num(totals.get("views", 0)), ""),
            ("Reach", _num(totals.get("reach", 0)), ""),
            ("Engagements", _num(totals.get("engagements", 0)), ""),
        ]
        _kpi_slide(prs, "Organic Performance", org_kpis)

        # Format breakdown table
        by_format = organic_summary.get("by_format", [])
        if by_format:
            headers = ["Format", "Posts", "Views", "Engagements", "ER"]
            rows = [[
                str(f.get("content_format", "Unknown")),
                _num(f.get("posts", 0)),
                _num(f.get("views", 0)),
                _num(f.get("engagements", 0)),
                f"{float(f.get('engagement_rate', 0)) * 100:.2f}%",
            ] for f in by_format[:6]]
            _data_table_slide(prs, "Performance by Content Format", headers, rows)

    # ── 5. Conversations Section ──
    if conversations and conversations.get("total_conversations"):
        _section_divider(prs, "Conversations", "Response time & messaging analysis")
        conv_kpis = [
            ("Conversations", _num(conversations.get("total_conversations", 0)), ""),
            ("Avg Response Time", conversations.get("avg_response_time_display", "N/A"), ""),
            ("Response Rate", f"{conversations.get('response_rate', 0):.0f}%", ""),
            ("Messages Received", _num(conversations.get("messages_received", 0)), ""),
        ]
        _kpi_slide(prs, "Messaging & Response Time", conv_kpis)

    # ── 6. Recommendations ──
    _bullets_slide(prs, "Recommended Next Steps", [
        "Convert top organic content into paid creative tests.",
        "Review high-spend campaigns with zero conversions before scaling budget.",
        "Refresh weak creatives where CTR is low or frequency is high.",
        "Keep UTM naming consistent for cross-channel attribution.",
        "Verify conversion events and tracking before any budget increase.",
    ])

    # ── 7. Closing ──
    _closing_slide(prs)

    prs.save(str(output_path))
    logger.info(f"Premium PPTX saved to {output_path}")
    return True


def generate_pptx_from_slides(slides_json_path: Path, output_path: Path) -> bool:
    """Legacy: Generate PPTX from AI slides JSON (kept for backward compatibility)."""
    if not PPTX_AVAILABLE:
        logger.warning("python-pptx is not installed. Skipping PPTX generation.")
        return False

    if not slides_json_path.exists():
        logger.warning(f"Slides JSON not found at {slides_json_path}")
        return False

    with slides_json_path.open("r", encoding="utf-8") as f:
        data = json.load(f)

    slides_data = data.get("slides", [])
    if not slides_data:
        logger.warning("No slides found in the JSON.")
        return False

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide_data in slides_data:
        slide_type = slide_data.get("type", "data")

        if slide_type == "title":
            _cover_slide(prs, slide_data.get("title", ""), slide_data.get("subtitle", ""))

        elif slide_type in ["summary", "action"]:
            bullets = slide_data.get("bullets", [])
            _bullets_slide(prs, slide_data.get("title", ""), bullets)

        elif slide_type == "data":
            insights = slide_data.get("insights", [])
            _bullets_slide(prs, slide_data.get("title", ""), insights)

        else:
            _bullets_slide(prs, slide_data.get("title", "Slide"), ["No content available."])

    _closing_slide(prs)
    prs.save(str(output_path))
    logger.info(f"PPTX Presentation saved to {output_path}")
    return True
